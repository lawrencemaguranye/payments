import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from urllib.parse import unquote

import pymupdf
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation

from postilion_agent.ingest.chunker import Section

HEADING_TAGS = {"h1": 1, "h2": 2, "h3": 3, "h4": 4, "h5": 5, "h6": 6}

_HHC_OBJECT_RE = re.compile(rb'<object[^>]*type="text/sitemap"[^>]*>(.*?)</object>', re.IGNORECASE | re.DOTALL)
_HHC_PARAM_RE = re.compile(rb'<param\s+name="([^"]+)"\s+value="([^"]*)"', re.IGNORECASE)

_MD_HEADER_RE = re.compile(r"^(#{1,6})\s+(.*)$")


class UnsupportedFileType(Exception):
    pass


def parse_pdf(path: Path) -> list[Section]:
    doc = pymupdf.open(path)
    try:
        toc = doc.get_toc(simple=True)
        if toc:
            return _pdf_sections_from_toc(doc, toc)
        return _pdf_sections_by_page(doc)
    finally:
        doc.close()


def _pdf_sections_from_toc(doc: "pymupdf.Document", toc: list[list]) -> list[Section]:
    sections: list[Section] = []
    stack: list[str] = []
    for i, (level, title, start_page) in enumerate(toc):
        end_page = toc[i + 1][2] if i + 1 < len(toc) else doc.page_count + 1
        stack = stack[: level - 1] + [title]
        parts = [
            doc[pno].get_text("text")
            for pno in range(start_page - 1, min(end_page - 1, doc.page_count))
            if pno >= 0
        ]
        text = "\n".join(parts).strip()
        if text:
            sections.append(Section(heading_path=list(stack), text=text, page=start_page))
    return sections


def _pdf_sections_by_page(doc: "pymupdf.Document") -> list[Section]:
    sections = []
    for pno in range(doc.page_count):
        text = doc[pno].get_text("text").strip()
        if text:
            sections.append(Section(heading_path=[], text=text, page=pno + 1))
    return sections


def parse_html_text(html: str, title_fallback: str) -> list[Section]:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    body = soup.body or soup

    sections: list[Section] = []
    stack: list[str] = []
    buffer: list[str] = []
    found_heading = False

    def flush() -> None:
        text = "\n".join(buffer).strip()
        if text:
            sections.append(Section(heading_path=list(stack), text=text))
        buffer.clear()

    for element in body.descendants:
        name = getattr(element, "name", None)
        if name in HEADING_TAGS:
            flush()
            title = element.get_text(" ", strip=True)
            if title:
                level = HEADING_TAGS[name]
                stack = stack[: level - 1] + [title]
                found_heading = True
        elif isinstance(element, NavigableString):
            parent_name = getattr(element.parent, "name", None)
            text = str(element).strip()
            if text and parent_name not in HEADING_TAGS:
                buffer.append(text)
    flush()

    if not found_heading and not sections:
        text = soup.get_text("\n", strip=True)
        if text:
            sections = [Section(heading_path=[title_fallback], text=text)]
    return sections


def parse_html(path: Path) -> list[Section]:
    html = path.read_text(encoding="utf-8", errors="ignore")
    return parse_html_text(html, title_fallback=path.stem)


def _find_7z() -> str:
    exe = shutil.which("7z") or shutil.which("7z.exe")
    if not exe:
        raise RuntimeError("7-Zip ('7z') not found on PATH - required to extract .chm files")
    return exe


def _parse_hhc_entries(raw: bytes) -> list[tuple[str, str]]:
    entries = []
    for match in _HHC_OBJECT_RE.finditer(raw):
        params = {}
        for name, value in _HHC_PARAM_RE.findall(match.group(1)):
            params[name.decode("utf-8", "ignore").lower()] = value.decode("utf-8", "ignore")
        title, local = params.get("name"), params.get("local")
        if title and local:
            entries.append((title, local))
    return entries


def parse_chm(path: Path) -> list[Section]:
    exe = _find_7z()
    extract_dir = Path(tempfile.mkdtemp(prefix="chm_", dir=".tmp")).resolve()
    try:
        subprocess.run(
            [exe, "x", str(path), f"-o{extract_dir}", "-y"],
            check=True,
            capture_output=True,
        )
        hhc_files = list(extract_dir.rglob("*.hhc"))
        entries = _parse_hhc_entries(hhc_files[0].read_bytes()) if hhc_files else []
        if entries:
            return _chm_sections_from_hhc(entries, extract_dir)
        return _chm_sections_from_html_files(extract_dir)
    finally:
        shutil.rmtree(extract_dir, ignore_errors=True)


def _chm_sections_from_hhc(entries: list[tuple[str, str]], extract_dir: Path) -> list[Section]:
    sections: list[Section] = []
    seen: set[Path] = set()
    for title, local in entries:
        file_part = unquote(local.split("#")[0])
        full_path = (extract_dir / file_part).resolve()
        if not full_path.exists() or full_path in seen:
            continue
        seen.add(full_path)
        html = full_path.read_text(encoding="utf-8", errors="ignore")
        for section in parse_html_text(html, title_fallback=title):
            heading = [title] + section.heading_path if section.heading_path else [title]
            sections.append(Section(heading_path=heading, text=section.text))
    return sections


def _chm_sections_from_html_files(extract_dir: Path) -> list[Section]:
    sections = []
    for html_path in sorted(extract_dir.rglob("*.htm*")):
        html = html_path.read_text(encoding="utf-8", errors="ignore")
        sections.extend(parse_html_text(html, title_fallback=html_path.stem))
    return sections


def parse_pptx(path: Path) -> list[Section]:
    presentation = Presentation(str(path))
    sections = []
    for i, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = None
        if title_shape is not None and title_shape.has_text_frame and title_shape.text_frame.text.strip():
            title = title_shape.text_frame.text.strip()
        title_id = getattr(title_shape, "shape_id", None) if title_shape is not None else None

        body_parts = []
        for shape in slide.shapes:
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)

        text = "\n".join(body_parts).strip()
        if not text and not title:
            continue
        heading = [title] if title else [f"Slide {i}"]
        sections.append(Section(heading_path=heading, text=text or title or "", slide=i))
    return sections


def parse_markdown(path: Path) -> list[Section]:
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    sections: list[Section] = []
    stack: list[str] = []
    buffer: list[str] = []

    def flush() -> None:
        text = "\n".join(buffer).strip()
        if text:
            sections.append(Section(heading_path=list(stack), text=text))
        buffer.clear()

    for line in lines:
        match = _MD_HEADER_RE.match(line)
        if match:
            flush()
            level = len(match.group(1))
            stack = stack[: level - 1] + [match.group(2).strip()]
        else:
            buffer.append(line)
    flush()
    return sections


_PARSERS = {
    ".pdf": parse_pdf,
    ".html": parse_html,
    ".htm": parse_html,
    ".chm": parse_chm,
    ".pptx": parse_pptx,
    ".md": parse_markdown,
    ".markdown": parse_markdown,
    ".txt": parse_markdown,
}


def parse_document(path: Path) -> list[Section]:
    parser = _PARSERS.get(path.suffix.lower())
    if parser is None:
        raise UnsupportedFileType(path.suffix)
    return parser(path)
